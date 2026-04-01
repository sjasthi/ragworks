# backend/rag_logic.py
import os
from openai import OpenAI, AuthenticationError
from dotenv import load_dotenv
from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
import chromadb
from chromadb.utils import embedding_functions
from datetime import datetime, timezone
import uuid

#SQL Integration
from database import get_connection

# Load environment variables
load_dotenv()

# OpenAI Client
def create_client(api_key):
    try:
        client = OpenAI(api_key=api_key)
        client.models.list()
        return client
    except AuthenticationError:
        print("Invalid API key.")
        return None

# Load API key from .env or ask
api_key = os.environ.get("OPENAI_KEY")
if not api_key:
    api_key = input("Enter your OpenAI API key: ").strip()

client = create_client(api_key)

# Persistent Chroma Setup
persist_dir = os.path.abspath("chroma_storage")
chroma_client = chromadb.PersistentClient(path=persist_dir)

embedding_function = embedding_functions.OpenAIEmbeddingFunction(
    api_key=api_key,
    model_name="text-embedding-3-small"
)

database = chroma_client.get_or_create_collection(
    name="RAG_Work_test",
    embedding_function=embedding_function
)

# checking in terminal for the eval_questions to see how many chunks are in the database before running the evaluation
print("CHUNKS:", database.count())

# --- Helper Functions ---

def normalize_path(path):
    return os.path.abspath(path)

def file_already_exists(file_path):
    file_path = normalize_path(file_path)
    results = database.get(where={"source": file_path}, limit=1)
    return len(results["ids"]) > 0

def chunk_text(text, chunk_size=800, overlap=150):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def now_iso():
    return datetime.now(timezone.utc).isoformat()

#MY SQL INTEGRATION Helper to log user queries and responses
def save_file_to_sql(file_path, file_type, uploaded_by=None):
    conn = get_connection()
    cursor = conn.cursor()

    filename = os.path.basename(file_path)

    cursor.execute("""
        INSERT INTO files (filename, file_type, file_path, uploaded_by)
        VALUES (%s, %s, %s, %s)
    """, (filename, file_type, file_path, uploaded_by))

    conn.commit()
    cursor.close()
    conn.close()

# --- File ingestion ---
def add_html(file_path, replace=False, uploaded_by=None):
    file_path = normalize_path(file_path)

    if file_already_exists(file_path):
        if not replace:
            return "File already uploaded."
        database.delete(where={"source": file_path})

    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    soup = BeautifulSoup(content, "html.parser")
    text = soup.get_text(separator=" ", strip=True)

    chunks = chunk_text(text)

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    upload_id = uuid.uuid4().hex[:8]
    uploaded_at = now_iso()

    database.add(
        documents=chunks,
        ids=[f"{file_name}_{upload_id}_chunk_{i}" for i in range(len(chunks))],
        metadatas=[{
            "source": file_path,
            "display_name": os.path.basename(file_path),
            "uploaded_at": uploaded_at,
            "upload_id": upload_id
        } for _ in chunks]
    )
    save_file_to_sql(file_path, "html", uploaded_by)
    return f"Added successfully. Collection size: {database.count()}"


#DocX is similar to HTML in that it needs a special library to extract the text.
def add_docx(file_path, replace=False, uploaded_by=None):
    file_path = normalize_path(file_path)

    if file_already_exists(file_path):
        if not replace:
            return "File already uploaded."
        database.delete(where={"source": file_path})

    doc = Document(file_path)
    text = "\n".join([p.text for p in doc.paragraphs])

    chunks = chunk_text(text)

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    upload_id = uuid.uuid4().hex[:8]
    uploaded_at = now_iso()

    database.add(
        documents=chunks,
        ids=[f"{file_name}_{upload_id}_chunk_{i}" for i in range(len(chunks))],
        metadatas=[{
            "source": file_path,
            "display_name": os.path.basename(file_path),
            "uploaded_at": uploaded_at,
            "upload_id": upload_id
        } for _ in chunks]
    )
    save_file_to_sql(file_path, "docx", uploaded_by)
    return f"Added successfully. Collection size: {database.count()}"

#TXT file is the simplest, just read and chunk without needing a special library
def add_txt(file_path, replace=False, uploaded_by=None):
    file_path = normalize_path(file_path)

    if file_already_exists(file_path):
        if not replace:
            return "File already uploaded."
        database.delete(where={"source": file_path})

    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    chunks = chunk_text(text)

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    upload_id = uuid.uuid4().hex[:8]
    uploaded_at = now_iso()

    database.add(
        documents=chunks,
        ids=[f"{file_name}_{upload_id}_chunk_{i}" for i in range(len(chunks))],
        metadatas=[{
            "source": file_path,
            "display_name": os.path.basename(file_path),
            "uploaded_at": uploaded_at,
            "upload_id": upload_id
        } for _ in chunks]
    )
    save_file_to_sql(file_path, "txt", uploaded_by)
    return f"Added successfully. Collection size: {database.count()}"

#PDF and PPTX are a bit more complex, so they have their own functions. 
def add_pdf(file_path, replace=False, uploaded_by=None):
    file_path = normalize_path(file_path)
    if file_already_exists(file_path):
        if not replace:
            return "File already uploaded"
        # admin replace: delete older chunks before adding new version
        database.delete(where={"source": file_path})

    reader = PdfReader(file_path)
    text = ""
    # because its a pdf, i will have to extract the text a bit differently using pypdf
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"

    chunks = chunk_text(text)
    file_name = os.path.splitext(os.path.basename(file_path))[0]
    upload_id = uuid.uuid4().hex[:8]
    uploaded_at = now_iso()

    database.add(
        documents=chunks,
        ids=[f"{file_name}_{upload_id}_chunk_{i}" for i in range(len(chunks))],
        metadatas=[{
            "source": file_path,
            "display_name": os.path.basename(file_path),
            "uploaded_at": uploaded_at,
            "upload_id": upload_id
        } for _ in chunks]
    )
    save_file_to_sql(file_path, "pdf", uploaded_by)
    return f"Added successfully. Collection size: {database.count()}"


# PPTX is similar to PDF in terms of needing a special library to extract text, but the structure is different so it gets its own function as well.
def add_pptx(file_path, replace=False, uploaded_by=None):
    file_path = normalize_path(file_path)
    if file_already_exists(file_path):
        if not replace:
            return "File already uploaded"
        # admin replace: delete older chunks before adding new version
        database.delete(where={"source": file_path})

    # same deal as pdf, need python-pptx to extract the text, using the offical guide as reference
    prs = Presentation(file_path)

    # basically making an array of texts from the presentation
    txt_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    txt_runs.append(run.text)

    text = " ".join(txt_runs)  # joining all of it into a single string so it can be passed onto chunks
    chunks = chunk_text(text)

    file_name = os.path.splitext(os.path.basename(file_path))[0]
    upload_id = uuid.uuid4().hex[:8]
    uploaded_at = now_iso()

    database.add(
        documents=chunks,
        ids=[f"{file_name}_{upload_id}_chunk_{i}" for i in range(len(chunks))],
        metadatas=[{
            "source": file_path,
            "display_name": os.path.basename(file_path),
            "uploaded_at": uploaded_at,
            "upload_id": upload_id
        } for _ in chunks]
    )
    save_file_to_sql(file_path, "pptx", uploaded_by)
    return f"Added successfully. Collection size: {database.count()}"


def add_file(file_path, file_type, replace=False, uploaded_by=None):
    file_type = file_type.lower()
    if file_type == "html":
        return add_html(file_path, replace=replace, uploaded_by=uploaded_by)
    elif file_type == "docx":
        return add_docx(file_path, replace=replace, uploaded_by=uploaded_by)
    elif file_type == "txt":
        return add_txt(file_path, replace=replace, uploaded_by=uploaded_by)
    elif file_type == "pdf":
        return add_pdf(file_path, replace=replace, uploaded_by=uploaded_by)
    elif file_type == "pptx":
        return add_pptx(file_path, replace=replace, uploaded_by=uploaded_by)
    else:
        return "Unsupported file type."

#Better version of get_contextto return the retrieved docs in a structured way so that run_rag can pass them to the evaluation function and we can see not just the answer but also what sources were retrieved for each question in the eval set. This is important for diagnosing whether low scores are due to retrieval failures or generation failures. Also added parameters for top_k, temperature, and top_p to run_rag so that we can test different settings in the eval.
def get_context(query, k=3, return_docs=False):
    results = database.query(
        query_texts=[query],
        n_results=k
    )

    documents = results["documents"][0]
    metadatas = results["metadatas"][0]

    retrieved_docs = []
    for doc, meta in zip(documents, metadatas):
        retrieved_docs.append({
            "content": doc,
            "source": meta.get("source", "unknown")
        })

    if return_docs:
        return retrieved_docs

    context_blocks = [
        f"[Source: {item['source']}]\n{item['content']}"
        for item in retrieved_docs
    ]
    return "\n\n".join(context_blocks)

#Original Run RAG function that just returns the generated answer without the retrieved sources. 
def run_rag(user_input):
    context = get_context(user_input)

    messages = [
        {
            "role": "system",
            "content": "Answer the question using ONLY the provided context. If not in context, say you do not know and do not include the source."
            "If the answer is found, include the source at the end on a seperate line in the format (Source: filename)."
        },
        {
            "role": "user",
            "content": f"Context:\n{context}\n\nQuestion: {user_input}"
        }
    ]

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
    )

    return response.choices[0].message.content

#Better Run Rag version that returns both the generated answer and the retrieved documents/sources in a structured way for better evaluation and diagnostics.
def run_rag_eval(user_input, top_k=3, temperature=0.0, top_p=1.0):
    retrieved_docs = get_context(user_input, k=top_k, return_docs=True)

    #Use only the basename file not the path in the source when constructing the context for the LLM making it cleaner for later output
    context = "\n\n".join(
        [f"[Source: {os.path.basename(doc['source'])}]\n{doc['content']}" 
        for doc in retrieved_docs]
    )

    messages = [
        {
            "role": "system",
            "content": "Answer the question using ONLY the provided context. If the answer is not in the context, say you do not know. When possible, mention the source used."
        },
        {
            "role": "user",
            "content": f"Context:\n{context}\n\nQuestion: {user_input}"
        }
    ]

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
        temperature=temperature,
        top_p=top_p
    )

    return {
        "answer": response.choices[0].message.content,
        "retrieved_docs": retrieved_docs
    }


#--- List and delete files Admin Helpers---

def list_documents():
    # Pull all metadatas; then dedupe by source
    data = database.get(include=["metadatas"])
    metadatas = data.get("metadatas", [])

    docs = {}
    for meta in metadatas:
        if not meta:
            continue
        source = meta.get("source")
        if not source:
            continue
        # Keep the "latest" uploaded_at if duplicates exist
        uploaded_at = meta.get("uploaded_at")
        display_name = meta.get("display_name") or os.path.basename(source)

        if source not in docs:
            docs[source] = {
                "source": source,
                "display_name": display_name,
                "uploaded_at": uploaded_at
            }
        else:
            # If you uploaded same source again, keep max uploaded_at
            if uploaded_at and (docs[source]["uploaded_at"] is None or uploaded_at > docs[source]["uploaded_at"]):
                docs[source]["uploaded_at"] = uploaded_at

    # return stable list sorted by uploaded_at desc (None last)
    return sorted(
        docs.values(),
        key=lambda d: (d["uploaded_at"] is None, d["uploaded_at"]),
        reverse=True
    )

def delete_document(source_path: str):
    source_path = normalize_path(source_path)
    # Delete all chunks with metadata source == source_path
    database.delete(where={"source": source_path})
    return {"deleted_source": source_path, "collection_size": database.count()}

def chroma_stats():
    docs = list_documents()
    return {
        "collection_name": "RAG_Work_test",
        "total_chunks": database.count(),
        "unique_documents": len(docs),
    }

# Admin upload newer version of file 
def add_file_admin(file_path, file_type, replace=False):
    file_path = normalize_path(file_path)
    if replace and file_already_exists(file_path):
        database.delete(where={"source": file_path})
    return add_file(file_path, file_type)

  


